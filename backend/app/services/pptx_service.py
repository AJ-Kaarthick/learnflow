from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text(file_path: Path) -> str:
    """
    Reads every slide of a .pptx file and returns its visible text, in
    slide order.

    Within a slide, shapes are walked in the order `slide.shapes`
    already lists them in (a slide's placeholder/body order is how
    python-pptx exposes shapes, and it matches how the slide reads
    top-to-bottom in the common case) -- that's what "preserving a
    logical reading order within each slide" means here, mirroring how
    docx_service preserves reading order via iter_inner_content().

    Mirrors pdf_service.extract_text / docx_service.extract_text's
    shape (one function, a `Path` in, a `str` out) and joins slides
    with blank lines between them for the same reason those two join
    pages/blocks that way: a visible break between distinct chunks of
    the document instead of one run-on wall of text.

    Deliberately narrow, per the milestone spec: only text already on
    the slide (titles, body placeholders, other text boxes, and text
    inside tables) is extracted. Speaker notes, animations, slide
    transitions, embedded audio/video, OCR from images, charts as
    structured data, and SmartArt semantics are all out of scope --
    a slide containing only an image or a chart simply contributes no
    text, the same way an image-only PDF page or a picture pasted into
    a DOCX contributes none in pdf_service / docx_service.
    """
    presentation = Presentation(str(file_path))

    slides_text: list[str] = []
    for slide in presentation.slides:
        slide_blocks = _slide_text_blocks(slide)
        if slide_blocks:
            slides_text.append("\n".join(slide_blocks))

    return "\n\n".join(slides_text).strip()


def _slide_text_blocks(slide) -> list[str]:
    """
    Returns the text-bearing blocks of a single slide, in shape order.
    A block is either a shape's plain text (a title, body, or other
    text box) or a table rendered one line per row -- same "cells
    joined with ' | '" convention docx_service uses for tables, so
    tabular content on a slide stays readable and grouped as plain
    text instead of losing all structure.
    """
    blocks: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_text = _table_to_text(shape.table)
            if table_text:
                blocks.append(table_text)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                blocks.append(text)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Grouped shapes don't have their own text_frame -- recurse
            # into the group so text inside a grouped title/box isn't
            # silently dropped.
            blocks.extend(_group_text_blocks(shape))
    return blocks


def _group_text_blocks(group_shape) -> list[str]:
    blocks: list[str] = []
    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_text = _table_to_text(shape.table)
            if table_text:
                blocks.append(table_text)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                blocks.append(text)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            blocks.extend(_group_text_blocks(shape))
    return blocks


def _table_to_text(table) -> str:
    """
    Renders a table as one line per row, cells separated by " | ".
    Rows that are entirely blank are dropped. Mirrors
    docx_service._table_to_text exactly, for the same reason.
    """
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)
