"""
V2.2 Milestone 3 — intelligent multi-document retrieval.

Covers what's new in retrieval_service.py and chat_service.py on top of
Milestone 2's multi-format support and the balanced multi-document
retrieval test_multi_document_chat.py already established: an adaptive
per-document chunk budget, duplicate/near-duplicate removal across the
merged result, comparison-question prompt selection, and informative
(rather than generic) failure messages when documents differ in
relevance. Doesn't re-prove anything test_multi_document_chat.py or
test_pdf_docx_pptx-format-specific suites already cover (e.g. that
balanced retrieval represents every document at all, or that PPTX
extraction works) — only what this milestone adds on top of that
foundation.
"""

import asyncio
import io
import random

import docx
from fastapi.testclient import TestClient
from pptx import Presentation
from reportlab.pdfgen import canvas

from app.db.database import SessionLocal
from app.db.models import Document
from app.main import app
from app.services.ai.base_provider import AIProvider
from app.services.ai.embedding_provider import EmbeddingProvider
from app.services.ai.embedding_provider_factory import get_embedding_provider
from app.services.ai.provider_factory import get_ai_provider
from app.services.chat_service import (
    NO_CONTEXT_ANSWER,
    answer_question,
    is_comparison_question,
)
from app.services.rag.embedding_service import index_document
from app.services.rag.retrieval_service import (
    MIN_ADAPTIVE_TOP_K,
    compute_adaptive_top_k,
    retrieve_relevant_chunks,
)

# ---------------------------------------------------------------------------
# Fakes and fixtures (local to this file — same convention every other test
# file in this suite follows; see docs/architecture.md's "Testing
# conventions").
# ---------------------------------------------------------------------------


class FakeAIProvider(AIProvider):
    """Returns a fixed answer and records the prompt it was called with."""

    def __init__(self, answer: str = "This is a fake grounded answer.") -> None:
        self._answer = answer
        self.last_prompt: str | None = None
        self.call_count = 0

    async def generate_text(self, prompt: str) -> str:
        self.last_prompt = prompt
        self.call_count += 1
        return self._answer


class TopicEmbeddingProvider(EmbeddingProvider):
    """
    Generalizes test_multi_document_chat.py's two-topic
    KeywordEmbeddingProvider to however many topics a given test needs.
    A text's vector is the (normalized) count of each topic keyword it
    contains — a query that's purely about one topic scores every other
    topic's chunks at exactly 0.0 (orthogonal), same sharp-relevance
    property the two-topic version relies on, just generalized to N
    topics instead of hardcoding two.
    """

    def __init__(self, topics: list[str]) -> None:
        self._topics = [topic.lower() for topic in topics]

    async def embed_document(self, text: str) -> list[float]:
        return self._vector_for(text)

    async def embed_query(self, text: str) -> list[float]:
        return self._vector_for(text)

    def _vector_for(self, text: str) -> list[float]:
        lowered = text.lower()
        counts = [float(lowered.count(topic)) for topic in self._topics]
        total = sum(counts)
        if total == 0:
            return [0.0] * len(self._topics)
        return [count / total for count in counts]


class ConstantEmbeddingProvider(EmbeddingProvider):
    """Every text maps to the same vector — retrieval always scores everything equally."""

    async def embed_document(self, text: str) -> list[float]:
        return [1.0]

    async def embed_query(self, text: str) -> list[float]:
        return [1.0]


_ANIMAL_WORDS = (
    "aardvark bison camel dingo elephant falcon gecko hedgehog iguana jackal "
    "kangaroo lemur mongoose narwhal octopus penguin quail raccoon salamander "
    "tortoise urchin vulture walrus xerus yak zebra badger cheetah dolphin "
    "egret ferret gazelle heron ibex jaguar koala llama macaw newt ocelot "
    "puma quokka rhino skunk tapir uakari viper wombat"
).split()


def _distinct_content_text(topic: str, paragraph_count: int = 60) -> str:
    """
    Generates `paragraph_count` short, lexically-varied paragraphs about
    `topic` — long enough, and different enough from each other, to
    reliably produce many distinct chunks under chunk_text() without
    tripping retrieval_service's near-duplicate detection. A template
    that only changed a trailing number was tried first and rejected:
    two chunking windows over that kind of text score as near-identical
    (retrieval correctly treats them as the same passage repeated),
    which is exactly right for real repeated boilerplate but wrong for
    a fixture that's supposed to simulate many genuinely different
    passages. Deterministic per topic (seeded on `topic` itself), so
    two calls with the same topic produce identical text and two
    different topics don't draw the same words at the same paragraph
    index.
    """
    rng = random.Random(topic)
    paragraphs = []
    for index in range(paragraph_count):
        chosen = rng.sample(_ANIMAL_WORDS, 6)
        paragraphs.append(
            f"{topic} note {index}: field observations mention {chosen[0]}, "
            f"{chosen[1]}, {chosen[2]}, {chosen[3]}, {chosen[4]} and "
            f"{chosen[5]} near the {topic} research site, recorded during "
            f"session {index} of the survey."
        )
    return " ".join(paragraphs)


def _make_indexed_document(
    db,
    embedding_provider: EmbeddingProvider,
    filename: str,
    text: str,
) -> Document:
    """
    Creates and indexes a Document directly against the DB (bypassing
    HTTP upload/index, same pattern test_chat.py's direct-service tests
    use) — the shortest path to "a document with real DocumentChunk
    rows scored by a specific embedding provider" for tests that only
    care about retrieval/chat behavior, not the upload pipeline itself.
    """
    document = Document(
        original_filename=filename,
        stored_filename=filename,
        status="ready",
        extracted_text=text,
    )
    db.add(document)
    db.commit()
    db.refresh(document)
    asyncio.run(index_document(document, db, embedding_provider))
    return document


def _make_test_pdf(text: str) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(50, 750, text)
    pdf.save()
    return buffer.getvalue()


def _make_test_docx(paragraphs: list[str]) -> bytes:
    document = docx.Document()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_test_pptx(slide_texts: list[str]) -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for index, text in enumerate(slide_texts):
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {index + 1}"
        slide.placeholders[1].text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _upload_and_index(client: TestClient, filename: str, content: bytes, content_type: str) -> str:
    upload_response = client.post(
        "/api/v1/documents/upload",
        files={"file": (filename, content, content_type)},
    )
    document_id = upload_response.json()["id"]
    client.post(f"/api/v1/documents/{document_id}/index")
    return document_id


# ---------------------------------------------------------------------------
# 1. Adaptive retrieval budget
# ---------------------------------------------------------------------------


def test_compute_adaptive_top_k_scales_down_as_document_count_grows():
    """
    Direct check of the tiered budget table itself: a single selected
    document gets the largest per-document budget, it shrinks as more
    documents are selected, and it never drops below the floor that
    guarantees every document still contributes useful evidence.
    """
    assert compute_adaptive_top_k(0) == 0
    assert compute_adaptive_top_k(1) == 8
    assert compute_adaptive_top_k(2) == 6
    assert compute_adaptive_top_k(3) == 5

    # 4+ documents: at the floor, and the floor itself is never
    # "one chunk per document" — this milestone explicitly rules that
    # out ("Do NOT use tiny fixed budgets like one chunk per document").
    for document_count in (4, 5, 7, 10, 25):
        assert compute_adaptive_top_k(document_count) == MIN_ADAPTIVE_TOP_K
        assert compute_adaptive_top_k(document_count) > 1

    # Monotonically non-increasing as document count grows — more
    # documents selected should never mean *more* budget per document.
    budgets = [compute_adaptive_top_k(n) for n in range(1, 11)]
    assert budgets == sorted(budgets, reverse=True)


def test_retrieval_budget_is_adaptive_across_document_counts():
    """
    End-to-end through retrieve_relevant_chunks: the per-document chunk
    budget actually used for a request (as reported by
    DocumentRetrievalSummary.retrieved_chunk_count) matches
    compute_adaptive_top_k(document_count) for that request — not one
    fixed number regardless of how many documents are selected. Every
    document here has far more indexed chunks available than any tier's
    budget, so the budget — not the document running out of content —
    is what's being measured.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        topics = ["Basalt", "Granite", "Limestone", "Marble", "Quartzite"]
        documents = [
            _make_indexed_document(
                db, embedding_provider, f"{topic.lower()}.pdf", _distinct_content_text(topic)
            )
            for topic in topics
        ]

        for document_count in (1, 2, 3, 5):
            selected = documents[:document_count]
            result = asyncio.run(
                retrieve_relevant_chunks(
                    document_ids=[document.id for document in selected],
                    query="Tell me about the rock samples",
                    db=db,
                    embedding_provider=embedding_provider,
                    top_k=None,
                )
            )
            expected_budget = compute_adaptive_top_k(document_count)
            assert len(result.document_summaries) == document_count
            for summary in result.document_summaries:
                assert summary.retrieved_chunk_count == expected_budget
    finally:
        db.close()


def test_explicit_top_k_overrides_adaptive_budget():
    """
    A caller that passes an explicit top_k always gets exactly that,
    regardless of how many documents are selected — the adaptive budget
    is only what fills in when a caller doesn't specify one (top_k=None),
    never a value that silently overrides an explicit request. This is
    what keeps a client-specified "give me N sources" request meaningful
    even though the chat schema's default is no longer a fixed number.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        doc_a = _make_indexed_document(
            db, embedding_provider, "a.pdf", _distinct_content_text("Alpha")
        )
        doc_b = _make_indexed_document(
            db, embedding_provider, "b.pdf", _distinct_content_text("Bravo")
        )

        result = asyncio.run(
            retrieve_relevant_chunks(
                document_ids=[doc_a.id, doc_b.id],
                query="anything",
                db=db,
                embedding_provider=embedding_provider,
                top_k=2,
            )
        )

        # Adaptive budget for 2 documents would be 6 — confirm the
        # explicit value of 2 was honored instead, for both documents.
        assert compute_adaptive_top_k(2) != 2
        for summary in result.document_summaries:
            assert summary.retrieved_chunk_count == 2
    finally:
        db.close()


# ---------------------------------------------------------------------------
# 2. Balanced retrieval
# ---------------------------------------------------------------------------


def test_balanced_retrieval_represents_every_document_despite_skewed_relevance():
    """
    A query purely about one of three topics scores the other two
    documents' chunks at exactly 0.0 by cosine similarity (orthogonal
    topic vectors) — the sharpest possible test of whether balanced,
    per-document retrieval still surfaces every selected document
    rather than letting the one clearly-relevant document dominate.
    """
    db = SessionLocal()
    try:
        embedding_provider = TopicEmbeddingProvider(["volcano", "glacier", "coral"])
        documents = [
            _make_indexed_document(
                db, embedding_provider, f"{topic}.pdf", _distinct_content_text(topic) + f" {topic} " * 20
            )
            for topic in ("volcano", "glacier", "coral")
        ]

        result = asyncio.run(
            retrieve_relevant_chunks(
                document_ids=[document.id for document in documents],
                query="Tell me about volcanoes",
                db=db,
                embedding_provider=embedding_provider,
                top_k=None,
            )
        )

        represented_document_ids = {scored.chunk.document_id for scored in result.chunks}
        assert represented_document_ids == {document.id for document in documents}

        # The two off-topic documents genuinely scored 0.0 — this isn't
        # balanced retrieval finding coincidental relevance, it's the
        # per-document guarantee doing its job despite zero relevance.
        by_id = {summary.document_id: summary for summary in result.document_summaries}
        assert by_id[documents[1].id].best_score == 0.0
        assert by_id[documents[2].id].best_score == 0.0
        assert by_id[documents[1].id].has_relevant_evidence is False
        assert by_id[documents[2].id].has_relevant_evidence is False
    finally:
        db.close()


def test_balanced_retrieval_is_unaffected_by_document_size():
    """
    A large document doesn't crowd out a small one: each document is
    retrieved from independently up to the adaptive budget, so a
    document with far more indexed chunks than another gets no more
    representation in the merged result than its budget allows, and a
    document with fewer chunks than the budget simply contributes all
    of what it has.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        large_document = _make_indexed_document(
            db, embedding_provider, "large.pdf", _distinct_content_text("Large", paragraph_count=120)
        )
        small_document = _make_indexed_document(
            db, embedding_provider, "small.pdf", "Small document with only one short paragraph of content."
        )

        result = asyncio.run(
            retrieve_relevant_chunks(
                document_ids=[large_document.id, small_document.id],
                query="anything",
                db=db,
                embedding_provider=embedding_provider,
                top_k=None,
            )
        )

        by_id = {summary.document_id: summary for summary in result.document_summaries}
        budget = compute_adaptive_top_k(2)
        assert by_id[large_document.id].retrieved_chunk_count == budget
        # The small document has fewer chunks than the budget, so it
        # contributes everything it has rather than being padded or
        # excluded — still represented in the merged result.
        assert 0 < by_id[small_document.id].retrieved_chunk_count <= budget
        assert any(
            scored.chunk.document_id == small_document.id for scored in result.chunks
        )
    finally:
        db.close()


# ---------------------------------------------------------------------------
# 3. Ranking behavior
# ---------------------------------------------------------------------------


def test_ranking_is_global_not_grouped_by_request_order():
    """
    The merged result is sorted by score across every selected document
    — not each document's chunks kept together in the order
    document_ids was given. Requesting the low-relevance document
    first and the high-relevance one second should still return the
    high-scoring chunks first.
    """
    db = SessionLocal()
    try:
        embedding_provider = TopicEmbeddingProvider(["oak", "pine"])
        oak_document = _make_indexed_document(
            db, embedding_provider, "oak.pdf", _distinct_content_text("oak") + " oak " * 20
        )
        pine_document = _make_indexed_document(
            db, embedding_provider, "pine.pdf", _distinct_content_text("pine") + " pine " * 20
        )

        # pine (irrelevant to the query) requested first, oak (relevant)
        # requested second.
        result = asyncio.run(
            retrieve_relevant_chunks(
                document_ids=[pine_document.id, oak_document.id],
                query="Tell me about oak trees",
                db=db,
                embedding_provider=embedding_provider,
                top_k=None,
            )
        )

        scores = [scored.score for scored in result.chunks]
        assert scores == sorted(scores, reverse=True)
        # The highest-scoring chunk overall must be from oak (score
        # 1.0), even though pine was listed first in document_ids.
        assert result.chunks[0].chunk.document_id == oak_document.id
        assert result.chunks[0].score == 1.0
    finally:
        db.close()


# ---------------------------------------------------------------------------
# 4. Duplicate / near-duplicate removal
# ---------------------------------------------------------------------------


def test_duplicate_boilerplate_chunk_is_removed_while_both_documents_stay_represented():
    """
    Two documents share a long, near-verbatim boilerplate passage (the
    kind of repeated disclaimer/notice this milestone's dedup exists
    to catch) in addition to their own unique content. After
    retrieval, only one copy of that shared passage should survive —
    but each document's own distinguishing content must still be
    present, proving dedup doesn't come at the cost of balanced
    representation.
    """
    db = SessionLocal()
    try:
        boilerplate = (
            "This course material is distributed under the Standard "
            "Academic Sharing Agreement version 4.2. Redistribution "
            "outside the enrolled cohort is not permitted without "
            "written consent from the instructor of record. All "
            "figures, diagrams and examples remain the intellectual "
            "property of the original author unless otherwise noted "
            "in an accompanying license file. This notice must be "
            "retained on every redistributed copy of this material, "
            "in whole or in part, regardless of format or medium. "
        ) * 2

        # Equal-length keyword prefixes so both documents' boilerplate
        # tail chunk lands at the exact same offset within the shared
        # text — otherwise chunk_text's word-boundary splitting shifts
        # the two documents' boilerplate chunks out of alignment and
        # they no longer read as near-duplicates of each other.
        prefix_alpha = "UNIQUE_TOPIC_ALPHA " * 8
        prefix_gamma = "UNIQUE_TOPIC_GAMMA " * 8
        assert len(prefix_alpha) == len(prefix_gamma)

        class AlphaGammaEmbeddingProvider(EmbeddingProvider):
            """
            The chunk containing a document's own unique keyword scores
            highest for a query mentioning both keywords; the pure
            boilerplate chunk (containing neither) scores 0.0 — so each
            document's unique chunk is always its own top (dedup-
            protected) chunk, and the boilerplate chunk is the one left
            exposed to duplicate removal.
            """

            async def embed_document(self, text: str) -> list[float]:
                return self._vector_for(text)

            async def embed_query(self, text: str) -> list[float]:
                return self._vector_for(text)

            @staticmethod
            def _vector_for(text: str) -> list[float]:
                lowered = text.lower()
                has_alpha = "alpha" in lowered
                has_gamma = "gamma" in lowered
                return [1.0 if has_alpha else 0.0, 1.0 if has_gamma else 0.0]

        embedding_provider = AlphaGammaEmbeddingProvider()
        doc_alpha = _make_indexed_document(
            db, embedding_provider, "alpha.pdf", prefix_alpha + " " + boilerplate
        )
        doc_gamma = _make_indexed_document(
            db, embedding_provider, "gamma.pdf", prefix_gamma + " " + boilerplate
        )

        result = asyncio.run(
            retrieve_relevant_chunks(
                document_ids=[doc_alpha.id, doc_gamma.id],
                query="Tell me about alpha and gamma",
                db=db,
                embedding_provider=embedding_provider,
                top_k=None,
            )
        )

        # Before dedup, both documents contributed their full 2 chunks
        # each (balanced retrieval saw everything).
        by_id = {summary.document_id: summary for summary in result.document_summaries}
        assert by_id[doc_alpha.id].retrieved_chunk_count == 2
        assert by_id[doc_gamma.id].retrieved_chunk_count == 2

        # After dedup: 4 candidate chunks in, but the two boilerplate
        # copies are near-identical, so only one should survive —
        # 3 chunks total, not 4.
        assert len(result.chunks) == 3

        # Both documents' unique, distinguishing chunk survived
        # (protected as each document's own top-scoring chunk).
        unique_chunk_texts = [
            scored.chunk.content for scored in result.chunks if "UNIQUE_TOPIC" in scored.chunk.content
        ]
        assert any("ALPHA" in text for text in unique_chunk_texts)
        assert any("GAMMA" in text for text in unique_chunk_texts)

        # Exactly one copy of the boilerplate-only passage remains.
        boilerplate_only_chunks = [
            scored for scored in result.chunks if "UNIQUE_TOPIC" not in scored.chunk.content
        ]
        assert len(boilerplate_only_chunks) == 1

        # Both documents are still represented in the final result —
        # dedup removed a redundant copy, not a whole document's
        # evidence.
        represented_document_ids = {scored.chunk.document_id for scored in result.chunks}
        assert represented_document_ids == {doc_alpha.id, doc_gamma.id}
    finally:
        db.close()


# ---------------------------------------------------------------------------
# 5. Comparison-question prompt selection
# ---------------------------------------------------------------------------


def test_is_comparison_question_detects_trigger_phrases():
    comparison_questions = [
        "Compare cats and dogs.",
        "What's the difference between TCP and UDP?",
        "What are the differences across these documents?",
        "What's the similarity between these two approaches?",
        "List the similarities and differences.",
        "Pros and cons of each approach?",
        "What are the advantages of option A?",
        "What are the disadvantages of option B?",
        "How do these compare across documents?",
        "Which document covers deadlocks?",
        "Approach A versus approach B.",
    ]
    for question in comparison_questions:
        assert is_comparison_question(question), question


def test_is_comparison_question_rejects_plain_questions():
    plain_questions = [
        "What is a deadlock?",
        "Summarize chapter 3.",
        "Explain how photosynthesis works.",
        "What does the author conclude?",
        "List the steps in the process.",
    ]
    for question in plain_questions:
        assert not is_comparison_question(question), question


def test_comparison_prompt_addition_present_only_for_multidoc_comparison_questions():
    """
    The dedicated comparison instructions (Milestone 3, requirement 5)
    should be added to the prompt only when both conditions hold: the
    question is comparison-style AND more than one document is
    selected. Neither condition alone is enough.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        doc_a = _make_indexed_document(db, embedding_provider, "a.pdf", _distinct_content_text("Alpha"))
        doc_b = _make_indexed_document(db, embedding_provider, "b.pdf", _distinct_content_text("Bravo"))

        marker = "Directly compare the documents"

        # (a) multi-document + comparison question -> present.
        comparison_provider = FakeAIProvider()
        asyncio.run(
            answer_question(
                documents=[doc_a, doc_b],
                question="Compare the two documents.",
                db=db,
                ai_provider=comparison_provider,
                embedding_provider=embedding_provider,
            )
        )
        assert marker in comparison_provider.last_prompt

        # (b) multi-document + non-comparison question -> absent.
        plain_provider = FakeAIProvider()
        asyncio.run(
            answer_question(
                documents=[doc_a, doc_b],
                question="Summarize the key points.",
                db=db,
                ai_provider=plain_provider,
                embedding_provider=embedding_provider,
            )
        )
        assert marker not in plain_provider.last_prompt

        # (c) single document + comparison-worded question -> absent
        # (nothing to compare against).
        single_doc_provider = FakeAIProvider()
        asyncio.run(
            answer_question(
                documents=[doc_a],
                question="Compare the arguments in this document.",
                db=db,
                ai_provider=single_doc_provider,
                embedding_provider=embedding_provider,
            )
        )
        assert marker not in single_doc_provider.last_prompt
    finally:
        db.close()


def test_comparison_prompt_still_includes_normal_grounding_rules():
    """
    A comparison question doesn't lose any of the normal grounding
    instructions — the comparison instructions are additive on top of
    the same rules every other question gets, not a replacement prompt.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        doc_a = _make_indexed_document(db, embedding_provider, "a.pdf", _distinct_content_text("Alpha"))
        doc_b = _make_indexed_document(db, embedding_provider, "b.pdf", _distinct_content_text("Bravo"))

        provider = FakeAIProvider()
        asyncio.run(
            answer_question(
                documents=[doc_a, doc_b],
                question="Compare the two documents.",
                db=db,
                ai_provider=provider,
                embedding_provider=embedding_provider,
            )
        )

        prompt = provider.last_prompt
        assert "ONLY the document" in prompt
        assert NO_CONTEXT_ANSWER in prompt
        assert '"Document 1"' in prompt
    finally:
        db.close()


# ---------------------------------------------------------------------------
# 6. Mixed PDF + DOCX + PPTX retrieval
# ---------------------------------------------------------------------------


def test_mixed_format_multi_document_retrieval_represents_every_document():
    """
    Balanced multi-document retrieval doesn't care what format a
    document was extracted from — a PDF, a DOCX, and a PPTX selected
    together should all be represented in a multi-document chat's
    sources, the same as if all three had been the same format.
    """
    app.dependency_overrides[get_ai_provider] = lambda: FakeAIProvider(
        "Drawing on all three selected documents."
    )
    app.dependency_overrides[get_embedding_provider] = lambda: TopicEmbeddingProvider(
        ["orbital", "tectonic", "cellular"]
    )
    client = TestClient(app)
    try:
        pdf_text = "orbital " * 150
        docx_paragraphs = ["tectonic " * 40] * 4
        pptx_slides = ["cellular " * 40] * 3

        pdf_id = _upload_and_index(client, "orbit.pdf", _make_test_pdf(pdf_text), "application/pdf")
        docx_id = _upload_and_index(
            client,
            "plates.docx",
            _make_test_docx(docx_paragraphs),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
        pptx_id = _upload_and_index(
            client,
            "biology.pptx",
            _make_test_pptx(pptx_slides),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        response = client.post(
            "/api/v1/documents/chat",
            json={
                "document_ids": [pdf_id, docx_id, pptx_id],
                "question": "Summarize what each document covers.",
            },
        )

        assert response.status_code == 200
        body = response.json()
        assert body["grounded"] is True
        source_document_ids = {source["document_id"] for source in body["sources"]}
        assert source_document_ids == {pdf_id, docx_id, pptx_id}

        # Grouped source metadata (requirement 7) covers every
        # requested document too, in request order.
        assert [group["document_id"] for group in body["sources_by_document"]] == [
            pdf_id,
            docx_id,
            pptx_id,
        ]
        for group in body["sources_by_document"]:
            assert len(group["sources"]) > 0
    finally:
        app.dependency_overrides.clear()


# ---------------------------------------------------------------------------
# 7. Informative failure handling
# ---------------------------------------------------------------------------


def test_informative_answer_when_only_some_documents_are_relevant():
    """
    When the model falls back to the generic NO_CONTEXT_ANSWER but
    retrieval's own per-document scoring shows some selected documents
    were relevant and others weren't, the answer should be upgraded to
    say which is which — not left as a bare "I couldn't answer."
    """
    db = SessionLocal()
    try:
        embedding_provider = TopicEmbeddingProvider(["seismology"])
        relevant_document = _make_indexed_document(
            db,
            embedding_provider,
            "seismology.pdf",
            _distinct_content_text("Seismology") + " seismology " * 30,
        )
        irrelevant_document = _make_indexed_document(
            db, embedding_provider, "cooking.pdf", _distinct_content_text("Cooking")
        )

        provider = FakeAIProvider(answer=NO_CONTEXT_ANSWER)
        result = asyncio.run(
            answer_question(
                documents=[relevant_document, irrelevant_document],
                question="Tell me about seismology.",
                db=db,
                ai_provider=provider,
                embedding_provider=embedding_provider,
            )
        )

        assert result.answer != NO_CONTEXT_ANSWER
        assert "seismology.pdf" in result.answer
        assert "cooking.pdf" in result.answer
        assert "1" in result.answer  # "1 of 2 selected documents"
        assert "2" in result.answer
        # Still grounded — the model was actually consulted with real
        # context, this is just a clearer restatement of the outcome.
        assert result.grounded is True
    finally:
        db.close()


def test_informative_answer_when_no_documents_are_relevant():
    """
    Same upgrade applies when *none* of the selected documents scored
    as relevant — still more informative than the bare generic
    sentence, since it names every document that was actually checked.
    """
    db = SessionLocal()
    try:
        embedding_provider = TopicEmbeddingProvider(["astronomy"])
        doc_a = _make_indexed_document(
            db, embedding_provider, "cooking.pdf", _distinct_content_text("Cooking")
        )
        doc_b = _make_indexed_document(
            db, embedding_provider, "gardening.pdf", _distinct_content_text("Gardening")
        )

        provider = FakeAIProvider(answer=NO_CONTEXT_ANSWER)
        result = asyncio.run(
            answer_question(
                documents=[doc_a, doc_b],
                question="Tell me about astronomy.",
                db=db,
                ai_provider=provider,
                embedding_provider=embedding_provider,
            )
        )

        assert result.answer != NO_CONTEXT_ANSWER
        assert "cooking.pdf" in result.answer
        assert "gardening.pdf" in result.answer
        assert "2" in result.answer
    finally:
        db.close()


def test_generic_answer_preserved_when_every_document_scores_as_relevant():
    """
    Regression lock for the deliberate design boundary: if retrieval's
    own scoring doesn't distinguish any document as *less* relevant
    than another (every document cleared the relevance threshold), the
    model's verbatim NO_CONTEXT_ANSWER is left untouched — there's no
    more informative story to tell than the model already gave, so
    this must never silently rewrite an answer the model actually
    produced.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        doc_a = _make_indexed_document(db, embedding_provider, "a.pdf", _distinct_content_text("Alpha"))
        doc_b = _make_indexed_document(db, embedding_provider, "b.pdf", _distinct_content_text("Bravo"))

        provider = FakeAIProvider(answer=NO_CONTEXT_ANSWER)
        result = asyncio.run(
            answer_question(
                documents=[doc_a, doc_b],
                question="What is the capital of France?",
                db=db,
                ai_provider=provider,
                embedding_provider=embedding_provider,
            )
        )

        assert result.answer == NO_CONTEXT_ANSWER
    finally:
        db.close()


def test_single_document_no_context_answer_is_never_rewritten():
    """
    The informative-failure upgrade is a multi-document-only behavior —
    single-document chat keeps returning the model's NO_CONTEXT_ANSWER
    verbatim, unchanged from before this milestone.
    """
    db = SessionLocal()
    try:
        embedding_provider = ConstantEmbeddingProvider()
        document = _make_indexed_document(
            db, embedding_provider, "solo.pdf", _distinct_content_text("Solo")
        )

        provider = FakeAIProvider(answer=NO_CONTEXT_ANSWER)
        result = asyncio.run(
            answer_question(
                documents=[document],
                question="What is the capital of France?",
                db=db,
                ai_provider=provider,
                embedding_provider=embedding_provider,
            )
        )

        assert result.answer == NO_CONTEXT_ANSWER
    finally:
        db.close()
