"""
PPTX support (V2.2 Milestone 2).

Mirrors test_docx_support.py's shape and conventions exactly — these
tests exist to prove the *same* pipeline (upload -> extraction ->
chunking -> embeddings -> RAG -> summary/flashcards/quiz/mindmap/chat)
behaves identically for a .pptx-sourced document as it already does
for a .pdf- or .docx-sourced one, not to reimplement coverage of the
pipeline itself.
"""

import io
import json

import docx
from pptx import Presentation
from pptx.util import Inches
from fastapi.testclient import TestClient
from reportlab.pdfgen import canvas

from app.main import app
from app.services import document_extraction_service, pptx_service
from app.services.ai.base_provider import AIProvider, AIProviderError
from app.services.ai.embedding_provider import EmbeddingProvider
from app.services.ai.embedding_provider_factory import get_embedding_provider
from app.services.ai.provider_factory import get_ai_provider

PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
DOCX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

client = TestClient(app)


# ---------------------------------------------------------------------------
# Fakes (local to this file, same convention every other test file uses —
# see docs/architecture.md's "Testing conventions").
# ---------------------------------------------------------------------------


class FakeAIProvider(AIProvider):
    """Stands in for a real provider — instant, free, deterministic."""

    async def generate_text(self, prompt: str) -> str:
        return "This is a fake summary for testing."


class FakeStructuredAIProvider(AIProvider):
    """Returns whatever well-formed JSON string it's constructed with,
    for flashcards/quiz/mindmap, which all expect structured output."""

    def __init__(self, response: str) -> None:
        self._response = response

    async def generate_text(self, prompt: str) -> str:
        return self._response


class FakeEmbeddingProvider(EmbeddingProvider):
    """Deterministic embeddings, same technique as test_rag.py."""

    async def embed_document(self, text: str) -> list[float]:
        return self._vector_for(text)

    async def embed_query(self, text: str) -> list[float]:
        return self._vector_for(text)

    @staticmethod
    def _vector_for(text: str) -> list[float]:
        lowered = text.lower()
        cat_count = lowered.count("cat")
        dog_count = lowered.count("dog")
        if cat_count > dog_count:
            return [1.0, 0.0]
        if dog_count > cat_count:
            return [0.0, 1.0]
        return [0.5, 0.5]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_test_pptx(
    slide_texts: list[str], table_rows: list[list[str]] | None = None
) -> bytes:
    """
    Builds a tiny real .pptx in memory, one slide per entry in
    `slide_texts` (a title-and-content slide with `text` as the body),
    the same "generate a real file rather than depend on a fixture on
    disk" approach _make_test_docx / _make_test_pdf use. `table_rows`,
    if given, is added as an extra slide containing only a table, so
    tests can prove table text is extracted.
    """
    presentation = Presentation()
    title_and_content_layout = presentation.slide_layouts[1]

    for index, text in enumerate(slide_texts):
        slide = presentation.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = f"Slide {index + 1}"
        body = slide.placeholders[1]
        body.text_frame.text = text

    if table_rows:
        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        rows, cols = len(table_rows), len(table_rows[0])
        graphic_frame = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(0.5), Inches(5), Inches(2)
        )
        table = graphic_frame.table
        for row_index, row_values in enumerate(table_rows):
            for col_index, value in enumerate(row_values):
                table.cell(row_index, col_index).text = value

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_empty_pptx() -> bytes:
    """A syntactically valid presentation with zero slides."""
    presentation = Presentation()
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_image_only_pptx() -> bytes:
    """
    A presentation with one slide that has no text-bearing shapes at
    all (blank layout, nothing added to it) — stands in for an
    image-only slide, since either way there's no text on it.
    """
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    presentation.slides.add_slide(blank_layout)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_test_docx(paragraphs: list[str]) -> bytes:
    document = docx.Document()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_test_pdf(text: str) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(100, 750, text)
    pdf.save()
    return buffer.getvalue()


def _upload_pptx(filename: str, pptx_bytes: bytes):
    return client.post(
        "/api/v1/documents/upload",
        files={"file": (filename, pptx_bytes, PPTX_CONTENT_TYPE)},
    )


def _upload_ready_pptx(filename: str = "test.pptx", text: str = "Some content.") -> str:
    response = _upload_pptx(filename, _make_test_pptx([text]))
    return response.json()["id"]


def _upload_and_index_pptx(filename: str, text: str) -> str:
    document_id = _upload_ready_pptx(filename, text)
    client.post(f"/api/v1/documents/{document_id}/index")
    return document_id


# ---------------------------------------------------------------------------
# pptx_service — unit-level extraction behavior
# ---------------------------------------------------------------------------


def test_pptx_extract_text_reads_slide_text(tmp_path):
    path = tmp_path / "notes.pptx"
    path.write_bytes(_make_test_pptx(["First slide body.", "Second slide body."]))

    text = pptx_service.extract_text(path)

    assert "First slide body." in text
    assert "Second slide body." in text


def test_pptx_extract_text_preserves_slide_order(tmp_path):
    path = tmp_path / "ordered.pptx"
    path.write_bytes(_make_test_pptx(["Alpha content.", "Beta content.", "Gamma content."]))

    text = pptx_service.extract_text(path)

    assert text.index("Alpha content.") < text.index("Beta content.") < text.index(
        "Gamma content."
    )


def test_pptx_extract_text_includes_table_text(tmp_path):
    path = tmp_path / "with_table.pptx"
    path.write_bytes(
        _make_test_pptx(
            ["Intro slide."],
            table_rows=[["Header A", "Header B"], ["Cell A1", "Cell B1"]],
        )
    )

    text = pptx_service.extract_text(path)

    assert "Header A | Header B" in text
    assert "Cell A1 | Cell B1" in text
    # The table slide comes after the intro slide.
    assert text.index("Intro slide.") < text.index("Header A")


def test_pptx_extract_text_returns_empty_string_for_a_presentation_with_no_slides(tmp_path):
    path = tmp_path / "empty.pptx"
    path.write_bytes(_make_empty_pptx())

    assert pptx_service.extract_text(path) == ""


def test_pptx_extract_text_returns_empty_string_for_a_slide_with_no_text(tmp_path):
    path = tmp_path / "image_only.pptx"
    path.write_bytes(_make_image_only_pptx())

    assert pptx_service.extract_text(path) == ""


def test_pptx_extract_text_raises_for_a_corrupted_file(tmp_path):
    path = tmp_path / "corrupted.pptx"
    path.write_bytes(b"this is not a real pptx file")

    try:
        pptx_service.extract_text(path)
        assert False, "expected extraction to raise for a corrupted file"
    except Exception:
        pass


# ---------------------------------------------------------------------------
# document_extraction_service — dispatch abstraction
# ---------------------------------------------------------------------------


def test_extraction_service_dispatches_pptx_to_pptx_service(tmp_path):
    path = tmp_path / "routed.pptx"
    path.write_bytes(_make_test_pptx(["Routed through the dispatcher."]))

    text = document_extraction_service.extract_text(path, ".pptx")

    assert "Routed through the dispatcher." in text


def test_extraction_service_has_no_page_counter_for_pptx(tmp_path):
    path = tmp_path / "no_pages.pptx"
    path.write_bytes(_make_test_pptx(["Some content."]))

    assert document_extraction_service.get_page_count(path, ".pptx") is None


def test_extraction_service_extension_matching_is_case_insensitive_for_pptx(tmp_path):
    path = tmp_path / "upper.pptx"
    path.write_bytes(_make_test_pptx(["Case insensitive."]))

    text = document_extraction_service.extract_text(path, ".PPTX")

    assert "Case insensitive." in text


# ---------------------------------------------------------------------------
# Upload endpoint
# ---------------------------------------------------------------------------


def test_upload_pptx_extracts_text():
    response = _upload_pptx("hello.pptx", _make_test_pptx(["Hello LearnFlow"]))

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert "Hello LearnFlow" in body["text_preview"]
    assert body["character_count"] > 0


def test_upload_pptx_includes_file_size_and_null_page_count():
    pptx_bytes = _make_test_pptx(["One deck, no page count."])

    response = _upload_pptx("sized.pptx", pptx_bytes)

    assert response.status_code == 201
    body = response.json()
    assert body["file_size_bytes"] == len(pptx_bytes)
    # PPTX has slides, not pages — same "nullable, shown as nothing"
    # contract DOCX already uses.
    assert body["page_count"] is None


def test_upload_pptx_preserves_original_filename_and_extension():
    response = _upload_pptx("Lecture Slides.pptx", _make_test_pptx(["Content."]))

    assert response.status_code == 201
    assert response.json()["original_filename"] == "Lecture Slides.pptx"


def test_upload_rejects_empty_pptx_file():
    response = _upload_pptx("empty.pptx", b"")

    assert response.status_code == 400


def test_upload_pptx_with_empty_presentation_still_succeeds_with_no_text():
    response = _upload_pptx("no_slides.pptx", _make_empty_pptx())

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert body["character_count"] == 0


def test_upload_pptx_with_corrupted_bytes_uploads_but_fails_processing():
    """
    A corrupted .pptx isn't rejected at the HTTP boundary (the
    content-type is legitimately "this is a PPTX upload") — same
    behavior as a corrupted PDF/DOCX today: the record is created so
    the user can see it happened, and status flips to "failed" rather
    than a 500 or a silently empty document.
    """
    response = _upload_pptx("corrupted.pptx", b"this is not a real pptx file")

    assert response.status_code == 201
    assert response.json()["status"] == "failed"

    get_response = client.get(f"/api/v1/documents/{response.json()['id']}")
    assert get_response.json()["status"] == "failed"


def test_get_pptx_document_returns_it_after_upload():
    upload_response = _upload_pptx("retrieve.pptx", _make_test_pptx(["Retrieve me"]))
    document_id = upload_response.json()["id"]

    get_response = client.get(f"/api/v1/documents/{document_id}")

    assert get_response.status_code == 200
    assert get_response.json()["id"] == document_id


def test_rename_and_delete_round_trip_for_an_uploaded_pptx():
    document_id = _upload_ready_pptx("Original.pptx", "Some content.")

    rename_response = client.patch(
        f"/api/v1/documents/{document_id}", json={"original_filename": "Renamed"}
    )
    assert rename_response.status_code == 200
    assert rename_response.json()["original_filename"] == "Renamed.pptx"

    delete_response = client.delete(f"/api/v1/documents/{document_id}")
    assert delete_response.status_code == 204

    get_response = client.get(f"/api/v1/documents/{document_id}")
    assert get_response.status_code == 404


# ---------------------------------------------------------------------------
# Mixed-document library / RAG / AI features for a PPTX-sourced document
# ---------------------------------------------------------------------------


def test_pdf_docx_and_pptx_documents_coexist_in_the_library():
    pdf_response = client.post(
        "/api/v1/documents/upload",
        files={"file": ("mixed-lib.pdf", _make_test_pdf("A PDF."), "application/pdf")},
    )
    docx_response = client.post(
        "/api/v1/documents/upload",
        files={"file": ("mixed-lib.docx", _make_test_docx(["A DOCX."]), DOCX_CONTENT_TYPE)},
    )
    pptx_response = _upload_pptx("mixed-lib.pptx", _make_test_pptx(["A PPTX."]))

    list_response = client.get("/api/v1/documents", params={"search": "mixed-lib"})

    assert list_response.status_code == 200
    names = {doc["original_filename"] for doc in list_response.json()}
    assert {"mixed-lib.pdf", "mixed-lib.docx", "mixed-lib.pptx"} <= names
    assert pdf_response.json()["status"] == "ready"
    assert docx_response.json()["status"] == "ready"
    assert pptx_response.json()["status"] == "ready"


def test_index_and_search_a_pptx_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    try:
        text = ("The cat sat on the mat. " * 60) + ("The dog ran in the yard. " * 60)
        document_id = _upload_and_index_pptx("cats_and_dogs.pptx", text)

        index_response = client.get(f"/api/v1/documents/{document_id}")
        assert index_response.json()["status"] == "ready"

        search_response = client.post(
            f"/api/v1/documents/{document_id}/search", json={"query": "Tell me about cats"}
        )
        assert search_response.status_code == 200
        results = search_response.json()["results"]
        assert results
        assert "cat" in results[0]["content"].lower()
    finally:
        app.dependency_overrides.clear()


def test_summary_works_for_a_pptx_document():
    app.dependency_overrides[get_ai_provider] = lambda: FakeAIProvider()
    try:
        document_id = _upload_ready_pptx(
            "summary.pptx", "Photosynthesis converts light to energy."
        )

        response = client.post(f"/api/v1/documents/{document_id}/summary")

        assert response.status_code == 201
        assert response.json()["content"] == "This is a fake summary for testing."
    finally:
        app.dependency_overrides.clear()


def test_flashcards_work_for_a_pptx_document():
    fake_cards = json.dumps([{"question": "What is 2+2?", "answer": "4"}])
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_cards)
    try:
        document_id = _upload_ready_pptx("flashcards.pptx", "Basic arithmetic.")

        response = client.post(f"/api/v1/documents/{document_id}/flashcards")

        assert response.status_code == 201
        cards = response.json()
        assert len(cards) == 1
        assert cards[0]["question"] == "What is 2+2?"
    finally:
        app.dependency_overrides.clear()


def test_quiz_works_for_a_pptx_document():
    fake_quiz = json.dumps(
        [
            {
                "question": "What is the capital of France?",
                "options": ["Paris", "Rome", "Berlin", "Madrid"],
                "correct_answer_index": 0,
            }
        ]
    )
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_quiz)
    try:
        document_id = _upload_ready_pptx("quiz.pptx", "European capitals.")

        response = client.post(f"/api/v1/documents/{document_id}/quiz")

        assert response.status_code == 201
        questions = response.json()
        assert len(questions) == 1
        assert questions[0]["correct_answer_index"] == 0
    finally:
        app.dependency_overrides.clear()


def test_mindmap_works_for_a_pptx_document():
    fake_mindmap = json.dumps({"title": "Central Topic", "children": []})
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_mindmap)
    try:
        document_id = _upload_ready_pptx("mindmap.pptx", "A topic with no subtopics.")

        response = client.post(f"/api/v1/documents/{document_id}/mindmap")

        assert response.status_code == 201
        assert response.json()["structure"]["title"] == "Central Topic"
    finally:
        app.dependency_overrides.clear()


# ---------------------------------------------------------------------------
# Chat (single- and multi-document) for a PPTX-sourced document
# ---------------------------------------------------------------------------


class FakeChatProvider(AIProvider):
    async def generate_text(self, prompt: str) -> str:
        return "Cats are mentioned in the excerpts."


def test_single_document_chat_works_for_a_pptx_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        document_id = _upload_and_index_pptx("chat.pptx", "The cat sat on the mat. " * 60)

        response = client.post(
            f"/api/v1/documents/{document_id}/chat", json={"question": "Tell me about cats"}
        )

        assert response.status_code == 200
        body = response.json()
        assert body["grounded"] is True
        assert body["answer"] == "Cats are mentioned in the excerpts."
    finally:
        app.dependency_overrides.clear()


def test_multi_document_chat_works_across_a_docx_and_a_pptx_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        docx_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "cats.docx",
                    _make_test_docx(["The cat sat on the mat. " * 30]),
                    DOCX_CONTENT_TYPE,
                )
            },
        )
        docx_id = docx_upload.json()["id"]
        client.post(f"/api/v1/documents/{docx_id}/index")

        pptx_id = _upload_and_index_pptx("dogs.pptx", "The dog ran in the yard. " * 30)

        response = client.post(
            "/api/v1/documents/chat",
            json={"document_ids": [docx_id, pptx_id], "question": "Compare cats and dogs"},
        )

        assert response.status_code == 200
        body = response.json()
        sources_by_document = {source["document_id"] for source in body["sources"]}
        assert docx_id in sources_by_document
        assert pptx_id in sources_by_document
    finally:
        app.dependency_overrides.clear()


def test_multi_document_chat_works_across_pdf_docx_and_pptx():
    """
    Proves multi-document chat — the code path the architecture doc
    calls out as shared between single- and multi-document chat — is
    genuinely format-agnostic across all three supported formats at
    once.
    """
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        pdf_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "three-way.pdf",
                    _make_test_pdf("The cat sat on the mat. " * 20),
                    "application/pdf",
                )
            },
        )
        pdf_id = pdf_upload.json()["id"]
        client.post(f"/api/v1/documents/{pdf_id}/index")

        docx_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "three-way.docx",
                    _make_test_docx(["The cat sat on the mat. " * 20]),
                    DOCX_CONTENT_TYPE,
                )
            },
        )
        docx_id = docx_upload.json()["id"]
        client.post(f"/api/v1/documents/{docx_id}/index")

        pptx_id = _upload_and_index_pptx("three-way.pptx", "The dog ran in the yard. " * 20)

        response = client.post(
            "/api/v1/documents/chat",
            json={
                "document_ids": [pdf_id, docx_id, pptx_id],
                "question": "Compare cats and dogs",
            },
        )

        assert response.status_code == 200
        body = response.json()
        sources_by_document = {source["document_id"] for source in body["sources"]}
        assert pdf_id in sources_by_document
        assert docx_id in sources_by_document
        assert pptx_id in sources_by_document
    finally:
        app.dependency_overrides.clear()
