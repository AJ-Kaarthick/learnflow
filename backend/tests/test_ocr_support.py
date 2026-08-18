"""
OCR support for scanned documents and images (V2.3 Milestone 1).

Mirrors test_docx_support.py / test_pptx_support.py's shape and
conventions exactly — these tests exist to prove the *same* pipeline
(upload -> extraction -> chunking -> embeddings -> RAG ->
summary/flashcards/quiz/mindmap/chat) behaves identically for an
image- or scanned-PDF-sourced document as it already does for a
text-layer PDF/DOCX/PPTX one, not to reimplement coverage of the
pipeline itself.

Also covers the two pieces this milestone actually adds on top of
that: ocr_service.py / the OCREngine abstraction (unit-level), and
document_extraction_service's PDF dispatch decision between
pdf_service (selectable text) and ocr_service (no selectable text) —
see the "PDF/OCR dispatch decision" section below.
"""

import io
import json

import docx
import pytest
from fastapi.testclient import TestClient
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from app.core.config import settings
from app.main import app
from app.services import document_extraction_service, ocr_service
from app.services.ai.base_provider import AIProvider
from app.services.ai.embedding_provider import EmbeddingProvider
from app.services.ai.embedding_provider_factory import get_embedding_provider
from app.services.ai.provider_factory import get_ai_provider
from app.services.ocr.ocr_engine_factory import get_ocr_engine
from app.services.ocr.pytesseract_engine import PytesseractOCREngine

PNG_CONTENT_TYPE = "image/png"
JPEG_CONTENT_TYPE = "image/jpeg"
DOCX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

client = TestClient(app)


# ---------------------------------------------------------------------------
# Fakes (local to this file, same convention every other test file uses —
# see docs/architecture.md's "Testing conventions").
# ---------------------------------------------------------------------------


class FakeAIProvider(AIProvider):
    """Stands in for a real provider — instant, free, deterministic."""

    async def generate_text(self, prompt: str) -> str:
        return "This is a fake summary for testing."


class FakeStructuredAIProvider(AIProvider):
    """Returns whatever well-formed JSON string it's constructed with,
    for flashcards/quiz/mindmap, which all expect structured output."""

    def __init__(self, response: str) -> None:
        self._response = response

    async def generate_text(self, prompt: str) -> str:
        return self._response


class FakeEmbeddingProvider(EmbeddingProvider):
    """Deterministic embeddings, same technique as test_rag.py."""

    async def embed_document(self, text: str) -> list[float]:
        return self._vector_for(text)

    async def embed_query(self, text: str) -> list[float]:
        return self._vector_for(text)

    @staticmethod
    def _vector_for(text: str) -> list[float]:
        lowered = text.lower()
        cat_count = lowered.count("cat")
        dog_count = lowered.count("dog")
        if cat_count > dog_count:
            return [1.0, 0.0]
        if dog_count > cat_count:
            return [0.0, 1.0]
        return [0.5, 0.5]


class FakeChatProvider(AIProvider):
    async def generate_text(self, prompt: str) -> str:
        return "Cats are mentioned in the excerpts."


# ---------------------------------------------------------------------------
# Helpers — building real images / scanned PDFs in memory (same
# "generate a real file rather than depend on a fixture on disk"
# approach _make_test_pdf / _make_test_docx / _make_test_pptx use
# elsewhere), and uploading them.
# ---------------------------------------------------------------------------

_FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"


def _render_text_image(text: str, font_size: int = 40) -> Image.Image:
    """
    Renders `text` onto a plain white canvas, large and high-contrast
    enough for Tesseract to read back reliably. The canvas is sized to
    the *actual measured* bounding box of the rendered text (via
    textbbox), not a guessed character-count width — a fixed
    px-per-character estimate is close enough for short words but
    reliably clips the last character or two off longer sentences,
    which silently truncates whatever a test then asserts got OCR'd
    back out, for a reason that has nothing to do with what the test
    is actually checking.
    """
    font = ImageFont.truetype(_FONT_PATH, font_size)
    probe = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    left, top, right, bottom = probe.textbbox((0, 0), text, font=font)

    padding = 40
    width = (right - left) + padding * 2
    height = (bottom - top) + padding * 2
    image = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((padding - left, padding - top), text, fill="black", font=font)
    return image


def _make_test_png(text: str) -> bytes:
    buffer = io.BytesIO()
    _render_text_image(text).save(buffer, format="PNG")
    return buffer.getvalue()


def _make_test_jpeg(text: str) -> bytes:
    buffer = io.BytesIO()
    _render_text_image(text).save(buffer, format="JPEG", quality=95)
    return buffer.getvalue()


def _make_blank_png() -> bytes:
    """A syntactically valid image with no text on it at all."""
    buffer = io.BytesIO()
    Image.new("RGB", (600, 160), color="white").save(buffer, format="PNG")
    return buffer.getvalue()


def _make_scanned_pdf(texts: list[str]) -> bytes:
    """
    Builds a "scanned" PDF: one page per string in `texts`, each
    containing only a rendered *image* of that text — no PDF text
    layer at all, the same way a real scanned or photographed
    document has none. This is exactly what pypdf's page.extract_text()
    is supposed to come back empty for, and what
    document_extraction_service's PDF dispatch is supposed to notice
    and fall back to OCR for.
    """
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    for text in texts:
        pdf.drawImage(ImageReader(_render_text_image(text)), 50, 600, width=400, height=100)
        pdf.showPage()
    pdf.save()
    return buffer.getvalue()


def _make_test_pdf(text: str) -> bytes:
    """A normal, text-layer PDF — used where a test needs a document
    the PDF/OCR dispatcher should *not* route to OCR."""
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(100, 750, text)
    pdf.save()
    return buffer.getvalue()


def _make_test_docx(paragraphs: list[str]) -> bytes:
    document = docx.Document()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_test_pptx(slide_texts: list[str]) -> bytes:
    presentation = Presentation()
    title_and_content_layout = presentation.slide_layouts[1]
    for index, text in enumerate(slide_texts):
        slide = presentation.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = f"Slide {index + 1}"
        slide.placeholders[1].text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _upload_png(filename: str, png_bytes: bytes):
    return client.post(
        "/api/v1/documents/upload",
        files={"file": (filename, png_bytes, PNG_CONTENT_TYPE)},
    )


def _upload_jpeg(filename: str, jpeg_bytes: bytes):
    return client.post(
        "/api/v1/documents/upload",
        files={"file": (filename, jpeg_bytes, JPEG_CONTENT_TYPE)},
    )


def _upload_ready_png(filename: str = "test.png", text: str = "Some content.") -> str:
    response = _upload_png(filename, _make_test_png(text))
    return response.json()["id"]


def _upload_and_index_png(filename: str, text: str) -> str:
    document_id = _upload_ready_png(filename, text)
    client.post(f"/api/v1/documents/{document_id}/index")
    return document_id


# ---------------------------------------------------------------------------
# OCREngine abstraction — unit-level
# ---------------------------------------------------------------------------


def test_pytesseract_engine_reads_text_from_an_image():
    engine = PytesseractOCREngine()

    text = engine.image_to_text(_render_text_image("Engine level text."))

    assert "Engine level text." in text


def test_ocr_engine_factory_returns_the_configured_engine():
    assert isinstance(get_ocr_engine(), PytesseractOCREngine)


def test_ocr_engine_factory_raises_for_an_unknown_engine():
    original = settings.ocr_engine
    settings.ocr_engine = "not-a-real-engine"
    try:
        with pytest.raises(ValueError):
            get_ocr_engine()
    finally:
        settings.ocr_engine = original


# ---------------------------------------------------------------------------
# ocr_service — unit-level extraction behavior
# ---------------------------------------------------------------------------


def test_ocr_service_extract_text_reads_a_png_image(tmp_path):
    path = tmp_path / "notes.png"
    path.write_bytes(_make_test_png("First line of notes."))

    text = ocr_service.extract_text(path)

    assert "First line of notes." in text


def test_ocr_service_extract_text_reads_a_jpeg_image(tmp_path):
    path = tmp_path / "notes.jpg"
    path.write_bytes(_make_test_jpeg("A JPEG photo of notes."))

    text = ocr_service.extract_text(path)

    assert "A JPEG photo of notes." in text


def test_ocr_service_extract_text_returns_empty_string_for_a_blank_image(tmp_path):
    path = tmp_path / "blank.png"
    path.write_bytes(_make_blank_png())

    assert ocr_service.extract_text(path) == ""


def test_ocr_service_extract_text_raises_for_a_corrupted_image_file(tmp_path):
    path = tmp_path / "corrupted.png"
    path.write_bytes(b"this is not a real image file")

    with pytest.raises(Exception):
        ocr_service.extract_text(path)


def test_ocr_service_extract_text_raises_for_a_truncated_image_file(tmp_path):
    """A file that starts out as a real PNG but is cut off partway
    through — the "half-uploaded" case a corrupted file could
    realistically look like, distinct from not-an-image-at-all."""
    real_bytes = _make_test_png("Will be truncated.")
    path = tmp_path / "truncated.png"
    path.write_bytes(real_bytes[: len(real_bytes) // 2])

    with pytest.raises(Exception):
        ocr_service.extract_text(path)


def test_ocr_service_extract_text_from_pdf_reads_a_scanned_pdf(tmp_path):
    path = tmp_path / "scanned.pdf"
    path.write_bytes(_make_scanned_pdf(["Scanned page content."]))

    text = ocr_service.extract_text_from_pdf(path)

    assert "Scanned page content." in text


def test_ocr_service_extract_text_from_pdf_joins_multiple_pages_with_blank_lines(tmp_path):
    path = tmp_path / "multi_scanned.pdf"
    path.write_bytes(_make_scanned_pdf(["Page one content.", "Page two content."]))

    text = ocr_service.extract_text_from_pdf(path)

    assert "Page one content." in text
    assert "Page two content." in text
    assert text.index("Page one content.") < text.index("Page two content.")
    # Mirrors pdf_service.extract_text's "blank line between pages" join.
    assert "\n\n" in text


# ---------------------------------------------------------------------------
# document_extraction_service — PDF/OCR dispatch decision
# ---------------------------------------------------------------------------


def test_extraction_service_dispatches_png_to_ocr_service(tmp_path):
    path = tmp_path / "routed.png"
    path.write_bytes(_make_test_png("Routed through the dispatcher."))

    text = document_extraction_service.extract_text(path, ".png")

    assert "Routed through the dispatcher." in text


def test_extraction_service_dispatches_jpg_to_ocr_service(tmp_path):
    path = tmp_path / "routed.jpg"
    path.write_bytes(_make_test_jpeg("Routed as a jpg."))

    text = document_extraction_service.extract_text(path, ".jpg")

    assert "Routed as a jpg." in text


def test_extraction_service_dispatches_jpeg_to_ocr_service(tmp_path):
    """A document whose extension is literally '.jpeg' (e.g. after a
    rename, which derives its extension from original_filename rather
    than the upload content-type map) still dispatches correctly."""
    path = tmp_path / "routed.jpeg"
    path.write_bytes(_make_test_jpeg("Routed as a jpeg."))

    text = document_extraction_service.extract_text(path, ".jpeg")

    assert "Routed as a jpeg." in text


def test_extraction_service_dispatches_a_scanned_pdf_to_ocr(tmp_path):
    path = tmp_path / "scanned.pdf"
    path.write_bytes(_make_scanned_pdf(["No selectable text here."]))

    text = document_extraction_service.extract_text(path, ".pdf")

    assert "No selectable text here." in text


def test_extraction_service_prefers_selectable_pdf_text_over_ocr(tmp_path, monkeypatch):
    """
    The core dispatch requirement: a PDF with a real text layer must
    use pdf_service, never OCR — proven here by making the OCR path
    raise if it's ever reached at all, not just by checking the
    output looks right.
    """
    path = tmp_path / "text.pdf"
    path.write_bytes(_make_test_pdf("Real selectable text."))

    def _fail_if_called(_path):
        raise AssertionError("OCR should not run for a PDF with selectable text")

    monkeypatch.setattr(ocr_service, "extract_text_from_pdf", _fail_if_called)

    text = document_extraction_service.extract_text(path, ".pdf")

    assert "Real selectable text." in text


def test_extraction_service_has_no_page_counter_for_images(tmp_path):
    path = tmp_path / "no_pages.png"
    path.write_bytes(_make_test_png("Some content."))

    assert document_extraction_service.get_page_count(path, ".png") is None
    assert document_extraction_service.get_page_count(path, ".jpg") is None
    assert document_extraction_service.get_page_count(path, ".jpeg") is None


def test_extraction_service_page_count_still_works_for_a_scanned_pdf(tmp_path):
    """OCR changes how a scanned PDF's *text* is obtained, never what
    "a page" means for a PDF — the page count still comes from the
    page tree, same as any other PDF."""
    path = tmp_path / "scanned_two_pages.pdf"
    path.write_bytes(_make_scanned_pdf(["Page one.", "Page two."]))

    assert document_extraction_service.get_page_count(path, ".pdf") == 2


def test_extraction_service_extension_matching_is_case_insensitive_for_images(tmp_path):
    path = tmp_path / "upper.PNG"
    path.write_bytes(_make_test_png("Case insensitive."))

    text = document_extraction_service.extract_text(path, ".PNG")

    assert "Case insensitive." in text


# ---------------------------------------------------------------------------
# Upload endpoint
# ---------------------------------------------------------------------------


def test_upload_png_extracts_text_via_ocr():
    response = _upload_png("hello.png", _make_test_png("Hello LearnFlow OCR"))

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert "Hello LearnFlow OCR" in body["text_preview"]
    assert body["character_count"] > 0


def test_upload_jpg_extracts_text_via_ocr():
    response = client.post(
        "/api/v1/documents/upload",
        files={"file": ("hello.jpg", _make_test_jpeg("Hello from a jpg"), JPEG_CONTENT_TYPE)},
    )

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert "Hello from a jpg" in body["text_preview"]


def test_upload_jpeg_extracts_text_via_ocr():
    response = _upload_jpeg("hello.jpeg", _make_test_jpeg("Hello from a jpeg"))

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert "Hello from a jpeg" in body["text_preview"]


def test_upload_scanned_pdf_extracts_text_via_ocr():
    response = client.post(
        "/api/v1/documents/upload",
        files={
            "file": (
                "scanned.pdf",
                _make_scanned_pdf(["This page is a scanned image."]),
                "application/pdf",
            )
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert "This page is a scanned image." in body["text_preview"]
    # Still a genuine PDF underneath — page count comes back same as
    # any other single-page PDF.
    assert body["page_count"] == 1


def test_upload_image_includes_file_size_and_null_page_count():
    png_bytes = _make_test_png("Sized image content.")

    response = _upload_png("sized.png", png_bytes)

    assert response.status_code == 201
    body = response.json()
    assert body["file_size_bytes"] == len(png_bytes)
    # Images have no concept of "pages" — nullable, same as DOCX/PPTX.
    assert body["page_count"] is None


def test_upload_preserves_original_filename_and_extension_for_an_image():
    response = _upload_png("Scanned Notes.png", _make_test_png("Content."))

    assert response.status_code == 201
    assert response.json()["original_filename"] == "Scanned Notes.png"


def test_upload_rejects_empty_png_file():
    response = _upload_png("empty.png", b"")

    assert response.status_code == 400


def test_upload_rejects_empty_jpg_file():
    response = client.post(
        "/api/v1/documents/upload",
        files={"file": ("empty.jpg", b"", JPEG_CONTENT_TYPE)},
    )

    assert response.status_code == 400


def test_upload_corrupted_image_uploads_but_fails_processing():
    """
    A corrupted image isn't rejected at the HTTP boundary (the
    content-type is legitimately "this is a PNG upload") — same
    behavior a corrupted PDF/DOCX/PPTX already has: the record is
    created so the user can see it happened, and status flips to
    "failed" rather than a 500 or a silently empty document.
    """
    response = _upload_png("corrupted.png", b"this is not a real png file")

    assert response.status_code == 201
    assert response.json()["status"] == "failed"

    get_response = client.get(f"/api/v1/documents/{response.json()['id']}")
    assert get_response.json()["status"] == "failed"


def test_upload_blank_image_succeeds_with_empty_extracted_text():
    """
    A syntactically valid image with no text on it isn't an error —
    same as a blank DOCX: OCR just comes back with nothing, the
    document is still "ready", and the empty result is visible in the
    metadata rather than looking like a failure.
    """
    response = _upload_png("blank.png", _make_blank_png())

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "ready"
    assert body["text_preview"] == ""
    assert body["character_count"] == 0


def test_upload_rejects_unsupported_file_type():
    response = client.post(
        "/api/v1/documents/upload",
        files={
            "file": (
                "sheet.xlsx",
                b"not supported",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )

    assert response.status_code == 400
    assert "PNG" in response.json()["detail"]
    assert "JPG" in response.json()["detail"]


def test_get_image_document_returns_it_after_upload():
    upload_response = _upload_png("retrieve.png", _make_test_png("Retrieve me"))
    document_id = upload_response.json()["id"]

    get_response = client.get(f"/api/v1/documents/{document_id}")

    assert get_response.status_code == 200
    assert get_response.json()["id"] == document_id


def test_pdf_docx_pptx_and_image_documents_coexist_in_the_library():
    pdf_response = client.post(
        "/api/v1/documents/upload",
        files={"file": ("mixed-ocr-lib.pdf", _make_test_pdf("A PDF."), "application/pdf")},
    )
    docx_response = client.post(
        "/api/v1/documents/upload",
        files={
            "file": (
                "mixed-ocr-lib.docx",
                _make_test_docx(["A DOCX."]),
                DOCX_CONTENT_TYPE,
            )
        },
    )
    pptx_response = client.post(
        "/api/v1/documents/upload",
        files={
            "file": (
                "mixed-ocr-lib.pptx",
                _make_test_pptx(["A PPTX."]),
                PPTX_CONTENT_TYPE,
            )
        },
    )
    png_response = _upload_png("mixed-ocr-lib.png", _make_test_png("A PNG."))

    list_response = client.get("/api/v1/documents", params={"search": "mixed-ocr-lib"})

    assert list_response.status_code == 200
    names = {doc["original_filename"] for doc in list_response.json()}
    assert {
        "mixed-ocr-lib.pdf",
        "mixed-ocr-lib.docx",
        "mixed-ocr-lib.pptx",
        "mixed-ocr-lib.png",
    } <= names
    for response in (pdf_response, docx_response, pptx_response, png_response):
        assert response.json()["status"] == "ready"


def test_rename_and_delete_round_trip_for_an_uploaded_image():
    document_id = _upload_ready_png("Original.png", "Some content.")

    rename_response = client.patch(
        f"/api/v1/documents/{document_id}", json={"original_filename": "Renamed"}
    )
    assert rename_response.status_code == 200
    assert rename_response.json()["original_filename"] == "Renamed.png"

    delete_response = client.delete(f"/api/v1/documents/{document_id}")
    assert delete_response.status_code == 204

    get_response = client.get(f"/api/v1/documents/{document_id}")
    assert get_response.status_code == 404


# ---------------------------------------------------------------------------
# RAG (chunking + embeddings + retrieval) for OCR-sourced documents
# ---------------------------------------------------------------------------


def test_index_and_search_an_image_sourced_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    try:
        document_id = _upload_and_index_png("cats.png", "The cat sat on the mat.")

        index_response = client.get(f"/api/v1/documents/{document_id}")
        assert index_response.json()["status"] == "ready"

        search_response = client.post(
            f"/api/v1/documents/{document_id}/search", json={"query": "Tell me about cats"}
        )
        assert search_response.status_code == 200
        results = search_response.json()["results"]
        assert results
        assert "cat" in results[0]["content"].lower()
    finally:
        app.dependency_overrides.clear()


def test_index_and_search_a_scanned_pdf_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    try:
        upload_response = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "scanned_dogs.pdf",
                    _make_scanned_pdf(["The dog ran in the yard."]),
                    "application/pdf",
                )
            },
        )
        document_id = upload_response.json()["id"]
        assert upload_response.json()["status"] == "ready"
        client.post(f"/api/v1/documents/{document_id}/index")

        search_response = client.post(
            f"/api/v1/documents/{document_id}/search", json={"query": "Tell me about dogs"}
        )
        assert search_response.status_code == 200
        results = search_response.json()["results"]
        assert results
        assert "dog" in results[0]["content"].lower()
    finally:
        app.dependency_overrides.clear()


# ---------------------------------------------------------------------------
# Summary / Flashcards / Quiz / Mind Map for an image-sourced document
# ---------------------------------------------------------------------------


def test_summary_works_for_an_image_document():
    app.dependency_overrides[get_ai_provider] = lambda: FakeAIProvider()
    try:
        document_id = _upload_ready_png(
            "summary.png", "Photosynthesis converts light to energy."
        )

        response = client.post(f"/api/v1/documents/{document_id}/summary")

        assert response.status_code == 201
        assert response.json()["content"] == "This is a fake summary for testing."
    finally:
        app.dependency_overrides.clear()


def test_flashcards_work_for_an_image_document():
    fake_cards = json.dumps([{"question": "What is 2+2?", "answer": "4"}])
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_cards)
    try:
        document_id = _upload_ready_png("flashcards.png", "Basic arithmetic.")

        response = client.post(f"/api/v1/documents/{document_id}/flashcards")

        assert response.status_code == 201
        cards = response.json()
        assert len(cards) == 1
        assert cards[0]["question"] == "What is 2+2?"
    finally:
        app.dependency_overrides.clear()


def test_quiz_works_for_an_image_document():
    fake_quiz = json.dumps(
        [
            {
                "question": "What is the capital of France?",
                "options": ["Paris", "Rome", "Berlin", "Madrid"],
                "correct_answer_index": 0,
            }
        ]
    )
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_quiz)
    try:
        document_id = _upload_ready_png("quiz.png", "European capitals.")

        response = client.post(f"/api/v1/documents/{document_id}/quiz")

        assert response.status_code == 201
        questions = response.json()
        assert len(questions) == 1
        assert questions[0]["correct_answer_index"] == 0
    finally:
        app.dependency_overrides.clear()


def test_mindmap_works_for_an_image_document():
    fake_mindmap = json.dumps({"title": "Central Topic", "children": []})
    app.dependency_overrides[get_ai_provider] = lambda: FakeStructuredAIProvider(fake_mindmap)
    try:
        document_id = _upload_ready_png("mindmap.png", "A topic with no subtopics.")

        response = client.post(f"/api/v1/documents/{document_id}/mindmap")

        assert response.status_code == 201
        assert response.json()["structure"]["title"] == "Central Topic"
    finally:
        app.dependency_overrides.clear()


# ---------------------------------------------------------------------------
# Chat (single-, multi-, and mixed-format) for OCR-sourced documents
# ---------------------------------------------------------------------------


def test_single_document_chat_works_for_an_image_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        document_id = _upload_and_index_png("chat.png", "The cat sat on the mat.")

        response = client.post(
            f"/api/v1/documents/{document_id}/chat", json={"question": "Tell me about cats"}
        )

        assert response.status_code == 200
        body = response.json()
        assert body["grounded"] is True
        assert body["answer"] == "Cats are mentioned in the excerpts."
    finally:
        app.dependency_overrides.clear()


def test_multi_document_chat_works_across_a_pdf_and_an_image_document():
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        pdf_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "cats.pdf",
                    _make_test_pdf("The cat sat on the mat. " * 20),
                    "application/pdf",
                )
            },
        )
        pdf_id = pdf_upload.json()["id"]
        client.post(f"/api/v1/documents/{pdf_id}/index")

        png_id = _upload_and_index_png("dogs.png", "The dog ran in the yard.")

        response = client.post(
            "/api/v1/documents/chat",
            json={"document_ids": [pdf_id, png_id], "question": "Compare cats and dogs"},
        )

        assert response.status_code == 200
        body = response.json()
        sources_by_document = {source["document_id"] for source in body["sources"]}
        assert pdf_id in sources_by_document
        assert png_id in sources_by_document
    finally:
        app.dependency_overrides.clear()


def test_multi_document_chat_works_across_pdf_docx_pptx_and_image():
    """
    Proves multi-document chat is genuinely format-agnostic across
    every format LearnFlow now supports at once, including one whose
    text came from OCR rather than a format-specific parser —
    retrieval and chat have no idea, and shouldn't need to.
    """
    app.dependency_overrides[get_embedding_provider] = lambda: FakeEmbeddingProvider()
    app.dependency_overrides[get_ai_provider] = lambda: FakeChatProvider()
    try:
        pdf_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "four-way.pdf",
                    _make_test_pdf("The cat sat on the mat. " * 20),
                    "application/pdf",
                )
            },
        )
        pdf_id = pdf_upload.json()["id"]
        client.post(f"/api/v1/documents/{pdf_id}/index")

        docx_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "four-way.docx",
                    _make_test_docx(["The cat sat on the mat. " * 20]),
                    DOCX_CONTENT_TYPE,
                )
            },
        )
        docx_id = docx_upload.json()["id"]
        client.post(f"/api/v1/documents/{docx_id}/index")

        pptx_upload = client.post(
            "/api/v1/documents/upload",
            files={
                "file": (
                    "four-way.pptx",
                    _make_test_pptx(["The dog ran in the yard. " * 20]),
                    PPTX_CONTENT_TYPE,
                )
            },
        )
        pptx_id = pptx_upload.json()["id"]
        client.post(f"/api/v1/documents/{pptx_id}/index")

        image_id = _upload_and_index_png("four-way.png", "The dog ran in the yard.")

        response = client.post(
            "/api/v1/documents/chat",
            json={
                "document_ids": [pdf_id, docx_id, pptx_id, image_id],
                "question": "Compare cats and dogs",
            },
        )

        assert response.status_code == 200
        body = response.json()
        sources_by_document = {source["document_id"] for source in body["sources"]}
        assert pdf_id in sources_by_document
        assert docx_id in sources_by_document
        assert pptx_id in sources_by_document
        assert image_id in sources_by_document
    finally:
        app.dependency_overrides.clear()
